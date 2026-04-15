"""
ChronoTable - Project Presentation Generator
Generates a professional PowerPoint presentation for the college timetable project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x24, 0x3B)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_ROSE = RGBColor(0xFB, 0x71, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if shape.has_text_frame:
        shape.text_frame.clear()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = "Segoe UI"
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_bullet_slide(slide, items, start_y=2.2, left=1.2, color=LIGHT_GRAY, size=17, spacing=0.48, icon="▸"):
    for i, item in enumerate(items):
        add_text(slide, left, start_y + i * spacing, 10, 0.45, f"{icon}  {item}", size=size, color=color)


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, 0, 0, 13.333, 0.12, ACCENT_BLUE)
add_shape(s, 0, 7.38, 13.333, 0.12, ACCENT_PURPLE)

add_text(s, 1, 1.5, 11, 0.6, "🎓", size=54, align=PP_ALIGN.CENTER)
add_text(s, 1, 2.3, 11, 0.9, "ChronoTable", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1, 3.1, 11, 0.6, "Intelligent College Timetable Generator", size=26, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_accent_line(s, 5.5, 3.9, 2.3, ACCENT_PURPLE)
add_text(s, 1, 4.2, 11, 0.5, "AI-powered scheduling with Genetic Algorithm & Constraint Satisfaction",
         size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 1, 5.5, 11, 0.4, "Project Presentation  •  2026", size=15, color=GRAY, align=PP_ALIGN.CENTER)

# Tech badges
techs = ["Next.js 14", "Express", "TypeScript", "PostgreSQL", "Prisma", "Genetic Algorithm"]
badge_start = 3.0
for i, tech in enumerate(techs):
    bx = badge_start + i * 1.25
    add_shape(s, bx, 6.1, 1.15, 0.35, BG_CARD, radius=True)
    add_text(s, bx, 6.12, 1.15, 0.35, tech, size=9, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: Problem Statement
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "Problem Statement", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_ROSE)

problems = [
    "Manual timetable creation takes days of tedious work each semester",
    "Frequent scheduling conflicts — teacher double-bookings & room clashes",
    "Difficult to balance teacher workloads fairly across the week",
    "Lab sessions require consecutive slots — hard to manage manually",
    "No centralized system for managing constraints & preferences",
    "Changes cascade into new conflicts, requiring full rework",
]
add_bullet_slide(s, problems, start_y=1.8, icon="✗")

add_shape(s, 1, 5.2, 11.3, 1.2, BG_CARD, radius=True)
add_text(s, 1.3, 5.35, 10.8, 0.9,
         "💡 Need: An intelligent, automated system that generates conflict-free, "
         "optimized timetables in seconds while respecting all institutional constraints.",
         size=16, color=ACCENT_GREEN)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: Proposed Solution
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "Proposed Solution", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_GREEN)

solutions = [
    ("🧬", "Hybrid Genetic Algorithm + CSP Engine", "Combines evolutionary optimization with constraint backtracking"),
    ("📊", "Full Academic Data Management", "CRUD for departments, teachers, subjects, rooms, batches, time slots"),
    ("⚙️", "18 Constraint Types", "Teacher, Subject, Room, Batch, and Global constraints with priority levels"),
    ("📅", "Multi-View Timetable Viewer", "View by Batch, Teacher, or Room with color-coded grids"),
    ("📈", "Analytics Dashboard", "Workload distribution, room utilization, AI-powered insights"),
    ("📥", "Multi-Format Export", "CSV, JSON, HTML, and Print-ready timetable output"),
]
for i, (icon, title, desc) in enumerate(solutions):
    y = 1.7 + i * 0.85
    add_text(s, 1, y, 0.5, 0.4, icon, size=22)
    add_text(s, 1.6, y, 5, 0.35, title, size=17, color=WHITE, bold=True)
    add_text(s, 1.6, y + 0.33, 10, 0.35, desc, size=14, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: Tech Stack
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "Technology Stack", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_BLUE)

stack = [
    ("Frontend", ["Next.js 14 (React)", "TypeScript", "Tailwind CSS", "ShadCN UI Components"], ACCENT_BLUE),
    ("Backend", ["Node.js + Express", "TypeScript", "JWT Authentication", "RESTful API"], ACCENT_PURPLE),
    ("Database", ["PostgreSQL / SQLite", "Prisma ORM", "11 Data Models", "Migration System"], ACCENT_GREEN),
    ("Algorithm", ["Genetic Algorithm", "CSP + Backtracking", "Conflict Scoring", "Population Evolution"], ACCENT_AMBER),
]
for i, (title, items, color) in enumerate(stack):
    x = 0.8 + i * 3.1
    add_shape(s, x, 1.7, 2.8, 4.5, BG_CARD, radius=True)
    add_accent_line(s, x + 0.3, 1.9, 2.2, color)
    add_text(s, x + 0.3, 2.1, 2.2, 0.4, title, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s, x + 0.3, 2.8 + j * 0.55, 2.2, 0.4, f"•  {item}", size=14, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: System Architecture
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "System Architecture", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_PURPLE)

# Frontend box
add_shape(s, 1, 1.8, 11.3, 1.4, BG_CARD, radius=True)
add_text(s, 1.3, 1.85, 3, 0.35, "Frontend — Next.js 14", size=16, color=ACCENT_BLUE, bold=True)
pages = ["Landing Page", "Login", "Dashboard", "Data Management", "Constraints", "Generate", "Timetable Viewer", "Analytics"]
for i, pg in enumerate(pages):
    add_shape(s, 1.3 + i * 1.35, 2.3, 1.25, 0.35, RGBColor(0x25, 0x30, 0x50), radius=True)
    add_text(s, 1.3 + i * 1.35, 2.32, 1.25, 0.35, pg, size=9, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 8, 1.85, 4, 0.35, "Tailwind CSS + ShadCN UI", size=12, color=GRAY, align=PP_ALIGN.RIGHT)

# Arrow
add_text(s, 6, 3.3, 1.5, 0.4, "▼  REST API", size=12, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Backend box
add_shape(s, 1, 3.8, 11.3, 1.4, BG_CARD, radius=True)
add_text(s, 1.3, 3.85, 3, 0.35, "Backend — Express + TypeScript", size=16, color=ACCENT_PURPLE, bold=True)
modules = ["Auth Service", "CRUD Routes", "Schedule Generation", "Analytics", "Constraint Validator", "Genetic Algorithm"]
for i, mod in enumerate(modules):
    add_shape(s, 1.3 + i * 1.85, 4.3, 1.75, 0.35, RGBColor(0x25, 0x30, 0x50), radius=True)
    add_text(s, 1.3 + i * 1.85, 4.32, 1.75, 0.35, mod, size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 8, 3.85, 4, 0.35, "Prisma ORM + JWT", size=12, color=GRAY, align=PP_ALIGN.RIGHT)

# Arrow
add_text(s, 6, 5.3, 1.5, 0.4, "▼  Prisma ORM", size=12, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)

# Database box
add_shape(s, 1, 5.8, 11.3, 1.2, BG_CARD, radius=True)
add_text(s, 1.3, 5.85, 3, 0.35, "Database — PostgreSQL", size=16, color=ACCENT_GREEN, bold=True)
tables = ["users", "departments", "semesters", "teachers", "subjects", "rooms", "batches", "time_slots", "constraints", "schedules"]
for i, tbl in enumerate(tables):
    add_shape(s, 1.3 + i * 1.1, 6.3, 1.0, 0.3, RGBColor(0x25, 0x30, 0x50), radius=True)
    add_text(s, 1.3 + i * 1.1, 6.3, 1.0, 0.3, tbl, size=8, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: Database Schema
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "Database Schema", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_GREEN)
add_text(s, 1, 1.5, 10, 0.4, "11 interconnected models managed via Prisma ORM with full relational integrity", size=15, color=GRAY)

models = [
    ("User", "id, email, password, name, role, teacherId", ACCENT_BLUE),
    ("Department", "id, name, code → teachers, subjects, batches", ACCENT_PURPLE),
    ("Teacher", "id, name, email, departmentId, maxLecturesDay", ACCENT_GREEN),
    ("Subject", "id, name, code, isLab, labDuration, lecturesPerWeek", ACCENT_AMBER),
    ("Room", "id, name, type (CLASSROOM|LAB), capacity", ACCENT_ROSE),
    ("Batch", "id, name, departmentId, semesterId, studentCount", ACCENT_BLUE),
    ("TimeSlot", "dayOfWeek, slotIndex, startTime, endTime, isBreak", ACCENT_PURPLE),
    ("Constraint", "type (18 types), priority, parameters (JSON)", ACCENT_GREEN),
    ("Schedule", "name, semesterId, status, fitnessScore", ACCENT_AMBER),
    ("ScheduleEntry", "scheduleId, subjectId, teacherId, roomId, batchId", ACCENT_ROSE),
]
for i, (name, fields, color) in enumerate(models):
    col = i % 2
    row = i // 2
    x = 1 + col * 6
    y = 2.1 + row * 1.0
    add_shape(s, x, y, 5.5, 0.85, BG_CARD, radius=True)
    add_text(s, x + 0.2, y + 0.05, 1.8, 0.35, name, size=15, color=color, bold=True)
    add_text(s, x + 0.2, y + 0.4, 5, 0.35, fields, size=11, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: Scheduling Algorithm
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 8, 0.7, "Scheduling Algorithm", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_AMBER)
add_text(s, 1, 1.5, 10, 0.4, "Hybrid Genetic Algorithm + Constraint Satisfaction Problem (CSP) Solver", size=16, color=GRAY)

steps = [
    ("1", "Build\nAllocations", "Map batch → subjects → teachers\nDetermine lecture counts", ACCENT_BLUE),
    ("2", "Initialize\nPopulation", "CSP backtracking + random\nN=50 chromosomes", ACCENT_PURPLE),
    ("3", "Evaluate\nFitness", "Score based on constraint\nviolations & penalties", ACCENT_GREEN),
    ("4", "Evolve\nPopulation", "Tournament selection →\nCrossover → Mutation", ACCENT_AMBER),
    ("5", "Elitism &\nConverge", "Preserve top 10%\n200 generations max", ACCENT_ROSE),
    ("6", "Output\nSchedule", "Return best chromosome\nas schedule entries", ACCENT_BLUE),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = 0.7 + i * 2.1
    add_shape(s, x, 2.2, 1.9, 2.8, BG_CARD, radius=True)
    add_shape(s, x + 0.65, 2.35, 0.6, 0.6, color, radius=True)
    add_text(s, x + 0.65, 2.35, 0.6, 0.6, num, size=22, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 3.15, 1.6, 0.6, title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 3.8, 1.6, 0.8, desc, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Constraint penalties table
add_text(s, 1, 5.4, 5, 0.4, "Constraint Penalties", size=18, color=WHITE, bold=True)
penalties = [
    ("Teacher/Room double-booking", "1000", "Hard"),
    ("Batch overlap", "1000", "Hard"),
    ("Lab in non-lab room", "200", "Medium"),
    ("Teacher overload/day", "50×extra", "Medium"),
    ("Workload imbalance", "5×σ", "Soft"),
]
for i, (constraint, penalty, severity) in enumerate(penalties):
    y = 5.9 + i * 0.3
    sev_color = ACCENT_ROSE if severity == "Hard" else ACCENT_AMBER if severity == "Medium" else ACCENT_GREEN
    add_text(s, 1.2, y, 5, 0.3, constraint, size=12, color=LIGHT_GRAY)
    add_text(s, 7, y, 1.5, 0.3, penalty, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 9, y, 1.5, 0.3, severity, size=12, color=sev_color, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: Constraint System
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "Constraint System", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_ROSE)
add_text(s, 1, 1.5, 10, 0.4, "18 constraint types across 5 categories with configurable priority levels", size=15, color=GRAY)

categories = [
    ("Teacher", ["Availability", "Max Lectures/Day", "Preferred Slots", "No Early Morning"], ACCENT_BLUE),
    ("Subject", ["Consecutive Labs", "Before Other", "Repetition Limit"], ACCENT_PURPLE),
    ("Room", ["Lab Only", "Capacity Check", "Department Specific"], ACCENT_GREEN),
    ("Batch", ["No Overlap", "Max Lectures", "Mandatory Break"], ACCENT_AMBER),
    ("Global", ["No Teacher Conflict", "No Room Conflict", "Balanced Workload", "No Heavy Consecutive"], ACCENT_ROSE),
]
for i, (cat, items, color) in enumerate(categories):
    x = 0.6 + i * 2.5
    add_shape(s, x, 2.1, 2.3, 3.8, BG_CARD, radius=True)
    add_accent_line(s, x + 0.3, 2.3, 1.7, color)
    add_text(s, x + 0.2, 2.5, 1.9, 0.4, cat, size=17, color=color, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s, x + 0.2, 3.2 + j * 0.45, 1.9, 0.35, f"•  {item}", size=12, color=LIGHT_GRAY)

# Priority levels
add_text(s, 1, 6.2, 3, 0.35, "Priority Levels:", size=15, color=WHITE, bold=True)
priorities = [("MANDATORY", ACCENT_ROSE), ("HIGH", ACCENT_AMBER), ("MEDIUM", ACCENT_BLUE), ("LOW", ACCENT_GREEN)]
for i, (p, c) in enumerate(priorities):
    add_shape(s, 4 + i * 2.2, 6.2, 2.0, 0.4, BG_CARD, radius=True)
    add_text(s, 4 + i * 2.2, 6.22, 2.0, 0.4, p, size=13, color=c, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: Features Overview
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "Key Features", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_BLUE)

features = [
    ("📊", "Data Management", "Full CRUD for all academic entities\nDepartments, Teachers, Subjects, Rooms, Batches", ACCENT_BLUE),
    ("🧬", "Smart Generation", "Sub-3-second generation for typical loads\nConfigurable algorithm parameters", ACCENT_PURPLE),
    ("📅", "Multi-View Timetable", "View by Batch, Teacher, or Room\nColor-coded grid with lab badges", ACCENT_GREEN),
    ("📥", "Export Options", "CSV, JSON, HTML, and Print formats\nBeautiful standalone HTML with dark theme", ACCENT_AMBER),
    ("📈", "Analytics Dashboard", "Teacher workload & room utilization charts\nAI-powered scheduling insights", ACCENT_ROSE),
    ("🔐", "Authentication", "JWT-based auth with Admin & Faculty roles\nProtected API endpoints", ACCENT_BLUE),
]
for i, (icon, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.7 + row * 2.7
    add_shape(s, x, y, 3.8, 2.4, BG_CARD, radius=True)
    add_accent_line(s, x + 0.3, y + 0.2, 1.5, color)
    add_text(s, x + 0.3, y + 0.4, 0.5, 0.5, icon, size=28)
    add_text(s, x + 0.9, y + 0.4, 2.5, 0.4, title, size=17, color=WHITE, bold=True)
    add_text(s, x + 0.3, y + 1.0, 3.2, 1.0, desc, size=13, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: API Endpoints
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "REST API Endpoints", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_PURPLE)

# Headers
add_shape(s, 1, 1.7, 11.3, 0.4, RGBColor(0x25, 0x30, 0x50))
add_text(s, 1.2, 1.72, 1.5, 0.35, "Method", size=13, color=ACCENT_BLUE, bold=True)
add_text(s, 3, 1.72, 4, 0.35, "Endpoint", size=13, color=ACCENT_BLUE, bold=True)
add_text(s, 7.5, 1.72, 4, 0.35, "Description", size=13, color=ACCENT_BLUE, bold=True)

endpoints = [
    ("POST", "/api/auth/login", "User authentication"),
    ("POST", "/api/auth/register", "User registration"),
    ("CRUD", "/api/departments", "Department management"),
    ("CRUD", "/api/teachers", "Teacher management"),
    ("CRUD", "/api/subjects", "Subject management (with teacher assignments)"),
    ("CRUD", "/api/rooms", "Room management"),
    ("CRUD", "/api/batches", "Batch management"),
    ("CRUD", "/api/time-slots", "Time slot management + bulk create"),
    ("CRUD", "/api/constraints", "Scheduling constraint management"),
    ("POST", "/api/schedules/generate", "Generate optimized timetable"),
    ("GET", "/api/schedules/:id", "Get schedule with entries (batch/teacher/room view)"),
    ("GET", "/api/analytics", "Analytics data & insights"),
]
for i, (method, endpoint, desc) in enumerate(endpoints):
    y = 2.2 + i * 0.38
    bg_color = BG_CARD if i % 2 == 0 else BG_DARK
    add_shape(s, 1, y, 11.3, 0.38, bg_color)
    m_color = ACCENT_GREEN if method == "GET" else ACCENT_AMBER if method == "POST" else ACCENT_PURPLE
    add_text(s, 1.2, y + 0.02, 1.5, 0.35, method, size=11, color=m_color, bold=True)
    add_text(s, 3, y + 0.02, 4, 0.35, endpoint, size=11, color=WHITE)
    add_text(s, 7.5, y + 0.02, 4, 0.35, desc, size=11, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: Results & Performance
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "Results & Performance", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_GREEN)

metrics = [
    ("⚡", "< 3 sec", "Generation Time", "For typical college workloads", ACCENT_BLUE),
    ("🎯", "99+", "Fitness Score", "Near-perfect constraint satisfaction", ACCENT_GREEN),
    ("🚫", "0", "Hard Conflicts", "Zero teacher/room/batch clashes", ACCENT_ROSE),
    ("📊", "18", "Constraint Types", "Comprehensive rule coverage", ACCENT_PURPLE),
]
for i, (icon, value, label, desc, color) in enumerate(metrics):
    x = 0.8 + i * 3.1
    add_shape(s, x, 1.7, 2.8, 2.5, BG_CARD, radius=True)
    add_text(s, x, 1.9, 2.8, 0.5, icon, size=32, align=PP_ALIGN.CENTER)
    add_text(s, x, 2.5, 2.8, 0.5, value, size=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, 3.1, 2.8, 0.35, label, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, 3.45, 2.8, 0.35, desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, 1, 4.6, 10, 0.4, "Key Achievements", size=20, color=WHITE, bold=True)
achievements = [
    "Conflict-free schedules generated consistently with genetic algorithm optimization",
    "Lab sessions correctly allocated to consecutive slots in appropriate lab rooms",
    "Balanced teacher workloads across the week using standard deviation minimization",
    "Real-time analytics with AI-powered insights for schedule quality assessment",
    "Multi-format export (CSV, JSON, HTML, Print) for institutional distribution",
    "Role-based access control with JWT authentication for Admin and Faculty users",
]
add_bullet_slide(s, achievements, start_y=5.1, size=14, spacing=0.38, icon="✓")

# ════════════════════════════════════════════════════════════════
# SLIDE 12: Future Scope
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 1, 0.5, 6, 0.7, "Future Scope", size=34, color=WHITE, bold=True)
add_accent_line(s, 1, 1.15, 2.5, ACCENT_AMBER)

future = [
    ("🤖", "AI Enhancement", "Deep learning-based preference prediction\nand adaptive constraint learning", ACCENT_BLUE),
    ("📱", "Mobile App", "React Native companion app for\nteachers to view schedules on-the-go", ACCENT_PURPLE),
    ("🔄", "Real-time Sync", "WebSocket-based live updates when\nschedules are modified or published", ACCENT_GREEN),
    ("🏫", "Multi-Campus", "Support for multi-campus institutions\nwith shared resource management", ACCENT_AMBER),
    ("📧", "Notifications", "Email/SMS alerts for schedule changes\nand automated reminders", ACCENT_ROSE),
    ("🔗", "LMS Integration", "Connect with Moodle, Canvas, and other\nLearning Management Systems", ACCENT_BLUE),
]
for i, (icon, title, desc, color) in enumerate(future):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.7 + row * 2.7
    add_shape(s, x, y, 3.8, 2.4, BG_CARD, radius=True)
    add_accent_line(s, x + 0.3, y + 0.2, 1.5, color)
    add_text(s, x + 0.3, y + 0.4, 0.5, 0.5, icon, size=28)
    add_text(s, x + 0.9, y + 0.4, 2.5, 0.4, title, size=17, color=WHITE, bold=True)
    add_text(s, x + 0.3, y + 1.0, 3.2, 1.0, desc, size=13, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 13: Thank You
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, 0, 0, 13.333, 0.12, ACCENT_BLUE)
add_shape(s, 0, 7.38, 13.333, 0.12, ACCENT_PURPLE)

add_text(s, 1, 2.0, 11, 0.8, "Thank You!", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(s, 5.5, 2.9, 2.3, ACCENT_BLUE)
add_text(s, 1, 3.2, 11, 0.5, "ChronoTable — Intelligent College Timetable Generator", size=22, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, 1, 4.2, 11, 0.5, "Questions & Discussion", size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_shape(s, 4, 5.2, 5.3, 1.2, BG_CARD, radius=True)
add_text(s, 4.3, 5.35, 4.7, 0.4, "🔗  github.com/chronotable", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, 4.3, 5.75, 4.7, 0.4, "📧  admin@college.edu", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Save ──
output_path = r"c:\Users\My Document\Desktop\project time table\ChronoTable_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
